#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
MAGI Obsidian Phase 2 – File Text Extractors

Supports: .md, .txt, .pdf, .docx
Optional:  .pptx, .xlsx (future)
"""

import hashlib
import logging
import os
import re
import shutil
import subprocess
from pathlib import Path
from typing import Dict, Optional

logger = logging.getLogger(__name__)

MAX_TEXT_CHARS = int(os.environ.get("MAGI_EXTRACT_MAX_CHARS", str(2_000_000)))  # 2M chars


def file_hash(path: Path) -> str:
    """Return a short sha256 hex of file contents (first 64KB for speed)."""
    h = hashlib.sha256()
    try:
        with open(path, "rb") as f:
            h.update(f.read(65536))
    except Exception:
        logging.getLogger(__name__).debug("silent-catch at %s:%s", __name__, 30, exc_info=True)
    return h.hexdigest()[:16]


def extract_text(path: Path) -> Dict:
    """Extract text from a file, dispatching by extension.

    Returns:
        {"success": True, "text": ..., "pages": ..., "method": ...}
        or {"success": False, "error": ...}
    """
    import os
    # MarkItDown path (feature-flagged, default OFF)
    if os.environ.get("MAGI_USE_MARKITDOWN", "0").strip() == "1":
        try:
            from skills.engine.document_reader import read_document
            r = read_document(str(path))
            if r.success and r.text:
                return {"success": True, "text": r.text, "pages": 1, "method": "markitdown"}
        except Exception:
            pass  # fall through to legacy

    ext = path.suffix.lower()
    try:
        if ext in (".md", ".txt", ".text", ".log", ".csv"):
            return _extract_plaintext(path)
        elif ext == ".pdf":
            return _extract_pdf(path)
        elif ext == ".docx":
            return _extract_docx(path)
        elif ext == ".doc":
            return _extract_doc_legacy(path)
        elif ext == ".pptx":
            return _extract_pptx(path)
        elif ext == ".xlsx":
            return _extract_xlsx(path)
        else:
            return {"success": False, "error": f"Unsupported extension: {ext}"}
    except Exception as e:
        return {"success": False, "error": f"{type(e).__name__}: {e}"}


# ── Plain text ──────────────────────────────────────────────────────

def _extract_plaintext(path: Path) -> Dict:
    for enc in ("utf-8", "big5", "cp950", "latin-1"):
        try:
            text = path.read_text(encoding=enc, errors="strict")
            return {"success": True, "text": text[:MAX_TEXT_CHARS], "pages": 1, "method": f"plaintext/{enc}"}
        except (UnicodeDecodeError, UnicodeError):
            continue
    # Last resort
    text = path.read_text(encoding="utf-8", errors="replace")
    return {"success": True, "text": text[:MAX_TEXT_CHARS], "pages": 1, "method": "plaintext/replace"}


# ── Multimodal extraction (RAG-Anything inspired) ─────────────────

def extract_text_multimodal(path: Path, *, case_number: str = "") -> Dict:
    """
    多模態文件提取 — 除了純文字外，也提取表格結構和圖片描述。
    使用 multimodal_parser 模組（靈感來自 RAG-Anything）。

    Returns:
        {"success": True, "text": structured_text, "pages": N,
         "method": "multimodal", "tables": N, "images": N,
         "parse_result": {...}}
    """
    if path.suffix.lower() != ".pdf":
        return extract_text(path)  # 非 PDF 走原始流程

    try:
        from skills.documents.multimodal_parser import parse_document
        result = parse_document(str(path), enable_llm_summary=True)

        if result.errors and not result.blocks:
            logger.debug("Multimodal parse failed for %s: %s, falling back", path, result.errors)
            return _extract_pdf(path)  # Fallback to standard extraction

        structured = result.structured_text
        if not structured or len(structured.strip()) < 20:
            return _extract_pdf(path)

        return {
            "success": True,
            "text": structured[:MAX_TEXT_CHARS],
            "pages": result.total_pages,
            "method": f"multimodal/{result.parser_used}",
            "tables": len(result.table_blocks),
            "images": len(result.image_blocks),
            "parse_result": result.to_dict(),
        }
    except ImportError:
        logger.debug("multimodal_parser not available, falling back to standard PDF extraction")
        return _extract_pdf(path)
    except Exception as e:
        logger.debug("Multimodal extraction failed for %s: %s, falling back", path, e)
        return _extract_pdf(path)


# ── PDF extraction (pdfplumber > PyMuPDF > PyPDF2 > OCR fallback) ──

_OCR_MIN_CHARS = 20          # below this threshold, text-layer is considered empty
_OCR_MAX_PAGES = 20          # cap OCR rendering to avoid stalling on huge PDFs
_OCR_DPI = 200               # resolution for page rendering


def _maybe_opendataloader_pdf(path: Path, current_text: str = "") -> Optional[Dict]:
    """Use OpenDataLoader PDF when the current extractor is weak."""
    if str(os.environ.get("MAGI_OPENDATALOADER_PDF_ENABLE", "auto")).strip().lower() in {"0", "false", "no", "off", "disabled"}:
        return None
    current = str(current_text or "").strip()
    if len(current) >= 500:
        return None
    try:
        from skills.engine.ocr import opendataloader_provider

        result = opendataloader_provider.run_pdf(str(path), task_type="legal")
    except Exception as exc:
        logger.debug("opendataloader provider failed for %s: %s", path, exc)
        return None
    text = (getattr(result, "corrected_text", "") or getattr(result, "raw_text", "") or "").strip()
    if not getattr(result, "success", False) or len(text) < max(_OCR_MIN_CHARS, len(current) + 100):
        return None
    return {
        "success": True,
        "text": text[:MAX_TEXT_CHARS],
        "pages": None,
        "method": getattr(result, "provider", "opendataloader_pdf"),
    }


def _extract_pdf(path: Path) -> Dict:
    page_count: Optional[int] = None

    # Try pdfplumber first (best for table-rich legal docs)
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(str(path)) as pdf:
            page_count = len(pdf.pages)
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
        combined = "\n\n".join(texts)
        if len(combined.strip()) >= _OCR_MIN_CHARS:
            odl = _maybe_opendataloader_pdf(path, combined)
            if odl:
                return odl
            return {"success": True, "text": combined[:MAX_TEXT_CHARS], "pages": page_count, "method": "pdfplumber"}
    except Exception as e:
        logger.debug("pdfplumber failed for %s: %s", path, e)

    # Fallback: PyMuPDF (fitz)
    try:
        import fitz
        doc = fitz.open(str(path))
        texts = []
        for page in doc:
            t = page.get_text()
            if t:
                texts.append(t)
        page_count = len(doc)
        doc.close()
        combined = "\n\n".join(texts)
        if len(combined.strip()) >= _OCR_MIN_CHARS:
            odl = _maybe_opendataloader_pdf(path, combined)
            if odl:
                return odl
            return {"success": True, "text": combined[:MAX_TEXT_CHARS], "pages": page_count, "method": "pymupdf"}
    except Exception as e:
        logger.debug("pymupdf failed for %s: %s", path, e)

    # Fallback: PyPDF2
    try:
        import PyPDF2
        with open(str(path), "rb") as f:
            reader = PyPDF2.PdfReader(f)
            page_count = len(reader.pages)
            texts = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
            combined = "\n\n".join(texts)
            if len(combined.strip()) >= _OCR_MIN_CHARS:
                odl = _maybe_opendataloader_pdf(path, combined)
                if odl:
                    return odl
                return {"success": True, "text": combined[:MAX_TEXT_CHARS], "pages": page_count, "method": "pypdf2"}
    except Exception as e:
        logger.debug("pypdf2 failed for %s: %s", path, e)

    # ── OCR fallback: render pages with PyMuPDF + tesseract ──────────
    odl = _maybe_opendataloader_pdf(path, "")
    if odl:
        return odl

    result = _extract_pdf_ocr(path, page_count)
    if result is not None:
        return result

    return {"success": False, "error": "All PDF extractors failed (including OCR)"}


def _extract_pdf_ocr(path: Path, page_count_hint: Optional[int] = None) -> Optional[Dict]:
    """Render PDF pages to images with PyMuPDF and OCR via tesseract.

    Feature flag MAGI_OBSIDIAN_OCR_CONSENSUS_ENABLE=0 (default off):
      - flag off  → legacy tesseract-only path (zero change)
      - flag on   → run_consensus() replaces single-page tesseract; failure fallback to legacy

    Returns an extraction dict on success, or None if OCR is unavailable.
    """
    # --- consensus feature flag -----------------------------------------------
    _consensus_enable = (
        str(os.environ.get("MAGI_OBSIDIAN_OCR_CONSENSUS_ENABLE", "0")).strip().lower()
        in {"1", "true", "yes", "on"}
    )

    try:
        import fitz
    except ImportError:
        logger.debug("PyMuPDF (fitz) not available; skipping OCR fallback for %s", path)
        return None

    # --- metrics writer (lazy, only used when consensus enabled) --------------
    def _write_metrics(page_num, img_hash, consensus_result, legacy_len):
        try:
            from api.platforms.runtime_dir import (
                metrics as _rt_metrics,
                atomic_append_jsonl as _rt_append,
            )
            import time as _time
            record = {
                "ts": _time.time(),
                "page_num": page_num,
                "img_hash": img_hash,
                "consensus_success": bool(
                    consensus_result and consensus_result.success
                ),
                "consensus_confidence": (
                    round(consensus_result.confidence, 4)
                    if consensus_result and consensus_result.success
                    else None
                ),
                "consensus_len": (
                    len(consensus_result.corrected_text)
                    if consensus_result and consensus_result.success
                    else 0
                ),
                "legacy_len": legacy_len,
                "mode": "enabled",
            }
            metrics_path = _rt_metrics("ocr") / "obsidian_ocr_consensus.jsonl"
            _rt_append(metrics_path, record, rotate_at=500, keep_tail=500)
        except Exception as e:
            logger.debug("obsidian: consensus metrics write failed: %s", e)

    # --- consensus runner (lazy import, lazy failure) -------------------------
    def _run_consensus_page(img_data_bytes):
        """Write PNG bytes to tmp file, run consensus, return result or None."""
        import tempfile
        tmp_path = None
        try:
            import importlib
            _ocr_mod = importlib.import_module("skills.engine.ocr.consensus")
            fd, tmp_path = tempfile.mkstemp(suffix=".png")
            try:
                os.write(fd, img_data_bytes)
            finally:
                os.close(fd)
            result = _ocr_mod.run_consensus(tmp_path, task_type="legal")
            if result and result.success:
                return result
            return None
        except Exception as e:
            logger.warning(
                "obsidian: consensus OCR failed, falling back to legacy: %s", e
            )
            return None
        finally:
            if tmp_path:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass

    # --- single-page legacy OCR (now through shared MAGI provider) ------------
    def _legacy_ocr_page(png_data):
        import tempfile
        tmp_path = None
        try:
            from skills.engine.ocr import tesseract_provider

            fd, tmp_path = tempfile.mkstemp(suffix=".png")
            try:
                os.write(fd, png_data)
            finally:
                os.close(fd)
            result = tesseract_provider.run(
                tmp_path,
                langs="chi_tra+eng",
                task_type="legal",
                timeout_sec=60,
            )
            if result and result.success:
                return (result.corrected_text or result.raw_text or "").strip()
            return ""
        except Exception:
            return ""
        finally:
            if tmp_path:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass

    try:
        import hashlib as _hashlib

        doc = fitz.open(str(path))
        page_count = len(doc)
        pages_to_ocr = min(page_count, _OCR_MAX_PAGES)
        ocr_texts = []

        for i in range(pages_to_ocr):
            page = doc[i]
            # Render page to PNG at the configured DPI
            pix = page.get_pixmap(dpi=_OCR_DPI)
            png_data = pix.tobytes("png")

            if not _consensus_enable:
                # --- flag off: pure legacy path (zero change) ---
                text = _legacy_ocr_page(png_data)
                if text:
                    ocr_texts.append(text)
            else:
                # --- flag on: consensus path, fallback to legacy ---
                img_hash = _hashlib.sha256(png_data).hexdigest()[:16]
                legacy_text = _legacy_ocr_page(png_data)
                consensus_result = _run_consensus_page(png_data)
                _write_metrics(i + 1, img_hash, consensus_result, len(legacy_text or ""))

                if consensus_result and consensus_result.success:
                    chosen = (
                        consensus_result.corrected_text
                        or consensus_result.selected_text
                        or ""
                    ).strip()
                    if chosen:
                        ocr_texts.append(chosen)
                        continue  # skip legacy append
                # consensus unavailable or empty → fall back
                if legacy_text:
                    ocr_texts.append(legacy_text)

        doc.close()

        if ocr_texts:
            combined = "\n\n".join(ocr_texts)[:MAX_TEXT_CHARS]
            mode = "consensus" if _consensus_enable else "tesseract"
            method = f"ocr/{mode} ({pages_to_ocr}/{page_count} pages)"
            return {"success": True, "text": combined, "pages": page_count, "method": method}

        logger.debug("OCR produced no text for %s", path)
        return None

    except Exception as e:
        logger.debug("OCR fallback failed for %s: %s", path, e)
        return None


# ── DOCX extraction ─────────────────────────────────────────────────

def _extract_docx(path: Path) -> Dict:
    import docx
    doc = docx.Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(paragraphs)[:MAX_TEXT_CHARS]
    return {"success": True, "text": text, "pages": max(1, len(paragraphs) // 40), "method": "python-docx"}


# ── Legacy .doc extraction (via textutil on macOS) ─────────────────

def _extract_doc_legacy(path: Path) -> Dict:
    """Extract text from legacy .doc files using macOS textutil."""
    if not shutil.which("textutil"):
        return {"success": False, "error": "textutil not available (macOS only)"}
    try:
        proc = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            capture_output=True, timeout=30,
        )
        if proc.returncode != 0:
            stderr = proc.stderr.decode("utf-8", errors="replace")[:200]
            return {"success": False, "error": f"textutil failed: {stderr}"}
        text = proc.stdout.decode("utf-8", errors="replace")
        if len(text.strip()) < 5:
            return {"success": False, "error": "textutil produced no text"}
        return {"success": True, "text": text[:MAX_TEXT_CHARS], "pages": 1, "method": "textutil"}
    except subprocess.TimeoutExpired:
        return {"success": False, "error": "textutil timed out"}
    except Exception as e:
        return {"success": False, "error": f"textutil: {e}"}


# ── PPTX extraction (optional) ──────────────────────────────────────

def _extract_pptx(path: Path) -> Dict:
    try:
        from pptx import Presentation
    except ImportError:
        return {"success": False, "error": "python-pptx not installed"}
    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
    text = "\n".join(texts)[:MAX_TEXT_CHARS]
    return {"success": True, "text": text, "pages": len(prs.slides), "method": "python-pptx"}


# ── XLSX extraction (optional) ──────────────────────────────────────

def _extract_xlsx(path: Path) -> Dict:
    try:
        import openpyxl
    except ImportError:
        return {"success": False, "error": "openpyxl not installed"}
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    texts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                texts.append("\t".join(cells))
    wb.close()
    text = "\n".join(texts)[:MAX_TEXT_CHARS]
    return {"success": True, "text": text, "pages": len(wb.sheetnames), "method": "openpyxl"}


# ── Metadata helpers ────────────────────────────────────────────────

_CASE_NUMBER_RE = re.compile(r"(\d{4})-(\d{4})")  # e.g. 2025-0002
_CASE_FOLDER_RE = re.compile(r"(\d{4}-\d{4})-(.+?)-(.*?)-(.*)")  # 2025-0002-Name-Phase-Charge


def parse_case_folder_name(folder_name: str) -> Dict:
    """Parse case metadata from folder name like '2025-0002-游秀鈴-一審-傷害致死'."""
    m = _CASE_FOLDER_RE.match(folder_name)
    if m:
        return {
            "case_number": m.group(1),
            "client_name": m.group(2),
            "phase": m.group(3),
            "charge": m.group(4),
        }
    return {}


def generate_note_title(source_path: Path, extracted_text: str) -> str:
    """Generate a short title for an extracted note."""
    stem = source_path.stem
    # Use first non-empty line of text as supplemental info
    for line in extracted_text.split("\n"):
        line = line.strip()
        if len(line) > 5:
            # Truncate long titles
            title_hint = line[:60].rstrip()
            return f"{stem}__{title_hint}" if len(stem) < 20 else stem
    return stem


SUPPORTED_EXTENSIONS = {".md", ".txt", ".text", ".log", ".csv", ".pdf", ".doc", ".docx", ".pptx", ".xlsx"}
